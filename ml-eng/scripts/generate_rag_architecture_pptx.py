#!/usr/bin/env python3
"""Generate OpenTrace RAG pipeline architecture PowerPoint (16:9)."""
from __future__ import annotations

import argparse
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from rag_architecture_content import (
    CORPUS_TABLE,
    DIAGRAM_SECTIONS,
    DIAGRAMS_PNG_DIR,
    DOC_VERSION,
    DOCS_DIR,
    INGEST_ARCHITECTURE_BULLETS,
    OUTPUT_PPTX,
    PURPOSE_BULLETS,
)
from render_architecture_diagrams import render_diagrams

# Widescreen 16:9
_SLIDE_W = Inches(13.333)
_SLIDE_H = Inches(7.5)

# Neutral dark-green / sand palette (not purple-gradient)
_BG = RGBColor(0x0F, 0x1A, 0x14)
_BG_ALT = RGBColor(0x16, 0x24, 0x1C)
_ACCENT = RGBColor(0x2F, 0x9E, 0x6B)
_ACCENT_WARM = RGBColor(0xC4, 0xA3, 0x5A)
_TEXT = RGBColor(0xF2, 0xF0, 0xE9)
_MUTED = RGBColor(0xA8, 0xB5, 0xAC)
_CARD = RGBColor(0x1A, 0x2B, 0x22)

_CAPTION_BY_STEM = {stem: caption for stem, _title, caption in DIAGRAM_SECTIONS}
_TITLE_BY_STEM = {stem: title for stem, title, _caption in DIAGRAM_SECTIONS}


def _set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, _SLIDE_W, _SLIDE_H)
    fill.fill.solid()
    fill.fill.fore_color.rgb = color
    fill.line.fill.background()
    # Send to back
    spTree = slide.shapes._spTree
    sp = fill._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def _add_text_box(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = _TEXT,
    align=PP_ALIGN.LEFT,
    font_name: str = "Calibri",
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name


def _add_bullets(
    slide,
    left,
    top,
    width,
    height,
    bullets: list[str],
    *,
    size: int = 16,
    color: RGBColor = _TEXT,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets[:6]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"•  {bullet}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Calibri"


def _accent_bar(slide, top=Inches(0)) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, _SLIDE_W, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _ACCENT
    bar.line.fill.background()


def _png(stem: str) -> Path:
    return DIAGRAMS_PNG_DIR / f"{stem}.png"


def _add_diagram(slide, stem: str, *, top=Inches(1.35), max_h=Inches(5.4)) -> None:
    path = _png(stem)
    if not path.is_file():
        _add_text_box(
            slide,
            Inches(0.8),
            top,
            Inches(11.5),
            Inches(1),
            f"[Missing diagram: {path.name}]",
            size=16,
            color=_ACCENT_WARM,
        )
        return
    # Fit image within margins; python-pptx preserves aspect if only width or height set.
    left = Inches(0.55)
    width = Inches(12.2)
    pic = slide.shapes.add_picture(str(path), left, top, width=width)
    # Cap height if oversized
    if pic.height > max_h:
        ratio = float(max_h) / float(pic.height)
        pic.height = max_h
        pic.width = int(pic.width * ratio)
        pic.left = int((_SLIDE_W - pic.width) / 2)


def _title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, _BG)
    _accent_bar(slide, Inches(2.35))
    _add_text_box(
        slide,
        Inches(0.9),
        Inches(2.55),
        Inches(11.5),
        Inches(1.2),
        "OpenTrace RAG Pipeline Architecture",
        size=36,
        bold=True,
        color=_TEXT,
    )
    _add_text_box(
        slide,
        Inches(0.9),
        Inches(3.7),
        Inches(11.5),
        Inches(0.5),
        f"Version {DOC_VERSION}  ·  Data pipelines → vector store & BigQuery → query-time LangGraph",
        size=16,
        color=_MUTED,
    )
    _add_text_box(
        slide,
        Inches(0.9),
        Inches(6.6),
        Inches(11.5),
        Inches(0.4),
        f"Generated {date.today().isoformat()}  ·  ml/rag",
        size=12,
        color=_MUTED,
    )


def _section_divider(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, _BG_ALT)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(2.9), Inches(1.2), Inches(0.12)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = _ACCENT
    bar.line.fill.background()
    _add_text_box(
        slide,
        Inches(0.9),
        Inches(3.2),
        Inches(11.5),
        Inches(1),
        title,
        size=32,
        bold=True,
        color=_TEXT,
    )
    _add_text_box(
        slide,
        Inches(0.9),
        Inches(4.2),
        Inches(11.5),
        Inches(0.6),
        subtitle,
        size=16,
        color=_MUTED,
    )


def _bullets_slide(prs: Presentation, title: str, bullets: list[str], *, footnote: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, _BG)
    _accent_bar(slide)
    _add_text_box(
        slide,
        Inches(0.7),
        Inches(0.35),
        Inches(12),
        Inches(0.7),
        title,
        size=26,
        bold=True,
        color=_TEXT,
    )
    _add_bullets(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.2), bullets, size=17)
    if footnote:
        _add_text_box(
            slide,
            Inches(0.8),
            Inches(6.7),
            Inches(11.5),
            Inches(0.4),
            footnote,
            size=11,
            color=_MUTED,
        )


def _diagram_slide(prs: Presentation, stem: str, *, title: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, _BG)
    _accent_bar(slide)
    headline = title or _TITLE_BY_STEM.get(stem, stem)
    caption = _CAPTION_BY_STEM.get(stem, "")
    _add_text_box(
        slide,
        Inches(0.55),
        Inches(0.28),
        Inches(12.2),
        Inches(0.55),
        headline,
        size=22,
        bold=True,
        color=_TEXT,
    )
    if caption:
        _add_text_box(
            slide,
            Inches(0.55),
            Inches(0.85),
            Inches(12.2),
            Inches(0.4),
            caption,
            size=12,
            color=_MUTED,
        )
    _add_diagram(slide, stem, top=Inches(1.3), max_h=Inches(5.7))


def _two_column_cards(
    prs: Presentation,
    title: str,
    left_title: str,
    left_bullets: list[str],
    right_title: str,
    right_bullets: list[str],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, _BG)
    _accent_bar(slide)
    _add_text_box(
        slide,
        Inches(0.7),
        Inches(0.35),
        Inches(12),
        Inches(0.6),
        title,
        size=26,
        bold=True,
        color=_TEXT,
    )
    for left, card_title, bullets in (
        (Inches(0.55), left_title, left_bullets),
        (Inches(6.85), right_title, right_bullets),
    ):
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.25), Inches(5.9), Inches(5.4)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = _CARD
        card.line.color.rgb = RGBColor(0x2A, 0x3F, 0x33)
        _add_text_box(
            slide,
            left + Inches(0.25),
            Inches(1.45),
            Inches(5.4),
            Inches(0.5),
            card_title,
            size=18,
            bold=True,
            color=_ACCENT,
        )
        _add_bullets(
            slide,
            left + Inches(0.25),
            Inches(2.1),
            Inches(5.4),
            Inches(4.2),
            bullets,
            size=14,
        )


def build_pptx(output: Path) -> Path:
    prs = Presentation()
    prs.slide_width = _SLIDE_W
    prs.slide_height = _SLIDE_H

    # Part A — Framing
    _title_slide(prs)
    _bullets_slide(
        prs,
        "Agenda",
        [
            "Data pipelines — how unstructured and structured sources land in stores",
            "Vector store — six Qdrant corpora, ingest ERD, hybrid-ready indexes",
            "BigQuery — staging_dev as the RAG SQL surface (not silver/gold at query time)",
            "Query-time RAG — control plane → vector + BQ legs → merge → generate → ACF",
            "Takeaways for product and engineering audiences",
        ],
    )
    _bullets_slide(
        prs,
        "Design principles",
        [str(b) for b in PURPOSE_BULLETS[:5]],
    )
    _diagram_slide(prs, "system_context")

    # Part B — Data pipelines
    _section_divider(
        prs,
        "Data pipelines",
        "Offline paths that fill Qdrant and BigQuery before any user query",
    )
    _two_column_cards(
        prs,
        "Two retrieval stores (co-equal)",
        "Unstructured → Qdrant",
        [
            "News, academic papers, policies, public reports, formation, OTA",
            "Drive / JSONL → preprocess → chunk → embed",
            "Dense E5 + optional BM25 sparse (hybrid RRF at query time)",
            "Payload indexes: geo, time, doc_kind",
        ],
        "Structured → BigQuery",
        [
            "Warehouse layers feed staging tables used by RAG",
            "Live SQL targets staging_dev only (never silver/gold)",
            "YAML catalog + measure ontology scope the reasoner",
            "Validated SELECT + LIMIT; templates / patterns / NL2SQL",
        ],
    )
    _bullets_slide(
        prs,
        "Vector ingest path",
        [str(b) for b in INGEST_ARCHITECTURE_BULLETS[:5]],
        footnote="See ingest ERD on next slide",
    )
    _diagram_slide(prs, "ingest_erd", title="Ingest and vector store ERD")
    _diagram_slide(prs, "six_corpora")
    _bullets_slide(
        prs,
        "Six Qdrant corpora (defaults)",
        [
            f"{row[0]} → {row[1]} — {row[2]}"
            for row in CORPUS_TABLE
            if row[0] != "Corpus"
        ],
    )
    _bullets_slide(
        prs,
        "BigQuery data path for RAG",
        [
            "Upstream bronze / silver / gold pipelines remain the warehouse source of truth",
            "RAG query surface is BQ_DATASET_SILVER (default staging_dev) — staging tables only",
            "Table contracts live under bq_tables_yaml_files/ (columns, joins, sample filters)",
            "Measure ontology maps food security, yield, prices, GDP, etc. to candidate tables",
            "Vector chunks may describe other layers; live SQL never queries them directly",
        ],
    )
    _diagram_slide(prs, "infra_erd", title="Infrastructure (services)")

    # Part C — Query-time RAG
    _section_divider(
        prs,
        "Query-time RAG",
        "LangGraph: control plane → vector leg → BQ leg → merge → generate",
    )
    _diagram_slide(prs, "entry_points")
    _diagram_slide(prs, "runtime_graph", title="Full LangGraph (E2E)")
    _diagram_slide(prs, "control_plane")
    _diagram_slide(prs, "measure_ontology")
    _diagram_slide(prs, "corpus_router")
    _diagram_slide(prs, "vector_retrieve")
    _diagram_slide(prs, "bq_retrieve_subflow")
    _diagram_slide(prs, "merge_rerank")
    _diagram_slide(prs, "generate_acf_subflow")
    _diagram_slide(prs, "plan_persona")

    # Part D — Close
    _diagram_slide(prs, "runtime_erd", title="Runtime domain ERD")
    _bullets_slide(
        prs,
        "Takeaways",
        [
            "Vector and BigQuery are peer retrieval legs that fuse at merge — not SQL-only RAG",
            "Offline ingest fills six Qdrant corpora; query-time SQL stays on staging_dev",
            "Control plane (enrich → decompose → ontology → task_mode) routes clarify vs full_rag",
            "ACF Path B scores cited OpenTrace evidence only (0–100); empty cites → no evidence",
            "LLM steps fail soft — heuristics / ontology fallbacks keep the graph from crashing",
        ],
    )
    _bullets_slide(
        prs,
        "Regenerate this deck",
        [
            "From ml-eng/:  python scripts/generate_rag_architecture_pptx.py",
            "Skip Mermaid re-render:  python scripts/generate_rag_architecture_pptx.py --skip-render",
            "Also available: generate_rag_architecture_pdf.py  ·  generate_rag_architecture_docx.py",
            "Output:  ml/rag/docs/OpenTrace-RAG-Pipeline-Architecture.pptx",
            f"Content version: {DOC_VERSION}",
        ],
    )

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    return output


def main() -> None:
    parser = argparse.ArgumentParser(description="Generate RAG pipeline architecture PPTX.")
    parser.add_argument(
        "--skip-render",
        action="store_true",
        help="Build from existing PNGs (skip mmdc).",
    )
    args = parser.parse_args()

    DOCS_DIR.mkdir(parents=True, exist_ok=True)
    missing = [s for s, _, _ in DIAGRAM_SECTIONS if not _png(s).is_file()]
    if missing:
        if args.skip_render:
            raise SystemExit(f"Missing PNGs (run without --skip-render): {missing}")
        render_diagrams()

    path = build_pptx(OUTPUT_PPTX)
    print(f"Wrote {path}")


if __name__ == "__main__":
    main()
